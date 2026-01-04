#!/usr/bin/env python3
"""
Claude Code 2.0 세미나 PPT 생성 스크립트
Design System: 테크/모던 스타일 (Dark Mode + Gradient)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Design System Colors
COLORS = {
    'primary': RGBColor(102, 126, 234),      # #667eea (Purple-Blue)
    'secondary': RGBColor(118, 75, 162),     # #764ba2 (Purple)
    'accent': RGBColor(240, 147, 251),       # #f093fb (Pink)
    'background': RGBColor(26, 26, 46),      # #1a1a2e (Dark Blue)
    'text_primary': RGBColor(255, 255, 255), # White
    'text_secondary': RGBColor(184, 184, 209), # #b8b8d1
    'success': RGBColor(78, 205, 196),       # #4ecdc4
    'warning': RGBColor(255, 230, 109),      # #ffe66d
    'error': RGBColor(255, 107, 107),        # #ff6b6b
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs

def add_dark_background(slide):
    """어두운 배경 추가"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['background']
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_text(slide, text, top, font_size=44, bold=True, color='text_primary'):
    """타이틀 텍스트 추가"""
    left = Inches(1)
    width = Inches(14)
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT

    return textbox

def add_bullet_text(slide, bullets, top, font_size=24):
    """불릿 포인트 추가"""
    left = Inches(1)
    width = Inches(14)
    height = Inches(4)

    textbox = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['text_secondary']
        p.font.name = 'Arial'
        p.space_before = Pt(12)

    return textbox

def add_code_block(slide, code, top, width=14, left=1):
    """코드 블록 추가"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(3))
    tf = textbox.text_frame
    tf.word_wrap = True

    # 배경 박스
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(30, 30, 50)

    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['success']
    p.font.name = 'Courier New'

    return textbox

def add_accent_box(slide, text, top, left=1, width=6, color='primary'):
    """강조 박스 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return shape

# ============================================
# SLIDE DEFINITIONS
# ============================================

def slide_01_title(prs):
    """슬라이드 1: 타이틀"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_dark_background(slide)

    add_title_text(slide, "Claude Code 2.0", 2.5, font_size=72, bold=True)
    add_title_text(slide, "완벽 활용 가이드", 3.8, font_size=56, bold=True, color='primary')
    add_title_text(slide, "AI 시대, 개발 생산성을 폭발시키는 방법", 5.5, font_size=28, bold=False, color='text_secondary')

    # 발표자 정보
    add_title_text(slide, "Woogi | Flutter Tech Lead", 7.5, font_size=20, bold=False, color='text_secondary')

def slide_02_agenda(prs):
    """슬라이드 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "오늘의 여정", 0.8, font_size=44)

    agenda_items = [
        "1. 왜 Claude Code인가?",
        "2. Context Engineering: 토큰의 경제학",
        "3. 병렬 Sub-agents: 시간을 사는 방법",
        "4. Skills & Commands: 팀 지식의 자산화",
        "5. MCP 통합: 개발 생태계 연결",
        "6. 실전 워크플로우 & AI 마인드셋",
        "7. 액션 아이템 & Q&A"
    ]
    add_bullet_text(slide, agenda_items, 2.2, font_size=28)

def slide_03_why_hook(prs):
    """슬라이드 3: 왜 Claude Code - Hook"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 1", 0.5, font_size=20, color='primary')
    add_title_text(slide, "개발 속도 2-3배 향상의 비밀", 1.2, font_size=44)

    # Before
    add_accent_box(slide, "Before (2024 초)", 2.8, left=1, width=6.5, color='error')
    before_items = ["코드 작성: 개발자 100%", "코드 리뷰: 수동", "문서화: \"나중에...\""]
    add_bullet_text(slide, before_items, 4.2, font_size=22)

    # After
    add_accent_box(slide, "After (현재)", 2.8, left=8.5, width=6.5, color='success')
    # After items positioned on right side
    textbox = slide.shapes.add_textbox(Inches(8.5), Inches(4.2), Inches(6.5), Inches(3))
    tf = textbox.text_frame
    for i, item in enumerate(["LLM 80% + 개발자 20%", "자동화 + 인간 검토", "코드와 동시 생성"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text_secondary']
        p.space_before = Pt(12)

def slide_04_why_opus(prs):
    """슬라이드 4: Opus 4.5 선택 이유"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Opus 4.5를 선택한 4가지 이유", 0.8, font_size=44)

    reasons = [
        "빠른 피드백 루프 → 진전이 체감됨 → 동기부여",
        "맥락 파악 탁월 → 의도를 정확히 감지",
        "협업적 대화 → 페어 프로그래밍 느낌",
        "최신 지식 (May 2025) → 최신 API 인지"
    ]

    for i, reason in enumerate(reasons):
        add_accent_box(slide, f"✅ {reason}", 2.2 + i * 1.5, left=1, width=14, color='primary')

def slide_05_context_concept(prs):
    """슬라이드 5: Context Window 개념"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 2: Context Engineering", 0.5, font_size=20, color='primary')
    add_title_text(slide, "Context Window = 200K 토큰 공간", 1.2, font_size=44)

    concept_text = """┌─────────────────────────────────────────────────┐
│            Context Window (200K tokens)          │
├─────────────────────────────────────────────────┤
│ User: "커피숍 랜딩페이지 만들어줘"                   │
│ Assistant: [tool_call: web_search]               │
│ Tool result: [검색 결과]    ← ~1.5K tokens 추가   │
│ Assistant: [tool_call: read_file]                │
│ Tool result: [파일 내용]    ← ~4K tokens 추가     │
│                                                  │
│ ⚠️ 모든 대화 + Tool 결과가 누적됨!                 │
│ ⚠️ LLM은 stateless → 매번 전체를 다시 읽음         │
└─────────────────────────────────────────────────┘"""
    add_code_block(slide, concept_text, 2.5)

def slide_06_context_rot(prs):
    """슬라이드 6: Context Rot 문제"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Context Rot: 성능 저하의 원인", 0.8, font_size=44)

    graph = """성능 ▲
100% │████████████████
     │                ████████
 75% │                        ████████
     │                                ████
 50% │                                    ████  ← Sweet Spot
     │                                        ████
 25% │                                            ▼ 급격한 저하
     │
  0% └──────┬──────┬──────┬──────┬──────┬──────▶ 사용률
           20%    40%    60%    80%   100%"""
    add_code_block(slide, graph, 2.2, width=12, left=2)

    add_accent_box(slide, "💡 50-60%에서 /compact 실행!", 6.5, left=4, width=8, color='warning')

def slide_07_statusline(prs):
    """슬라이드 7: Status Line 설정"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Status Line으로 실시간 모니터링", 0.8, font_size=44)

    add_title_text(slide, "가장 쉬운 방법:", 2.0, font_size=24, color='success')
    code1 = "/statusline show context percentage and model name"
    add_code_block(slide, code1, 2.6, width=12, left=1)

    add_title_text(slide, "결과 화면:", 4.2, font_size=24, color='success')
    result = "📁 ~/my-project 🌿 main 🤖 Opus 4.5 🟡 Context: 45%"
    add_code_block(slide, result, 4.8, width=12, left=1)

    add_accent_box(slide, "🟢 <50% | 🟡 50-70% | 🔴 >70%", 6.5, left=3, width=10, color='primary')

def slide_08_context_action(prs):
    """슬라이드 8: Context 관리 전략"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Context 관리 액션 플랜", 0.8, font_size=44)

    actions = [
        "50% 도달 → /compact 정리 고려",
        "70% 도달 → 반드시 /compact 또는 새 대화",
        "새 작업 시작 → /clear 로 깨끗하게",
        "이전 작업 이어가기 → /handoff → /clear"
    ]

    for i, action in enumerate(actions):
        color = ['success', 'warning', 'primary', 'secondary'][i]
        add_accent_box(slide, action, 2.2 + i * 1.5, left=1, width=14, color=color)

def slide_09_subagent_concept(prs):
    """슬라이드 9: Sub-agent 개념"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 3: 병렬 Sub-agents", 0.5, font_size=20, color='primary')
    add_title_text(slide, "독립 Context의 분신들", 1.2, font_size=44)

    diagram = """┌─────────────────────────────────────────────────┐
│              Main Agent (Opus 4.5)               │
│              Context: 200K tokens                │
└────────────────────┬────────────────────────────┘
                     │ spawn
       ┌─────────────┼─────────────┐
       ▼             ▼             ▼
  ┌─────────┐   ┌─────────┐   ┌─────────┐
  │ Agent 1 │   │ Agent 2 │   │ Agent 3 │
  │ 200K    │   │ 200K    │   │ 200K    │
  │ 독립 ctx │   │ 독립 ctx │   │ 독립 ctx │
  └────┬────┘   └────┬────┘   └────┬────┘
       └─────────────┴─────────────┘
                     ▼
              결과 종합/검증"""
    add_code_block(slide, diagram, 2.3)

def slide_10_comparison(prs):
    """슬라이드 10: 순차 vs 병렬 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "순차 vs 병렬: 시간 62% 단축!", 0.8, font_size=44)

    # 왼쪽: 순차
    add_accent_box(slide, "순차 처리", 2.2, left=1, width=6, color='error')
    seq = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(6), Inches(3))
    tf = seq.text_frame
    for text in ["토큰: ~70K", "시간: 8분", "Context 사용: 35%"]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_secondary']
        p.space_before = Pt(8)

    # 오른쪽: 병렬
    add_accent_box(slide, "4 병렬 처리", 2.2, left=9, width=6, color='success')
    para = slide.shapes.add_textbox(Inches(9), Inches(3.5), Inches(6), Inches(3))
    tf = para.text_frame
    for text in ["토큰: ~175K (2.5x)", "시간: 3분 (62%↓)", "Context 사용: 17.5%"]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_secondary']
        p.space_before = Pt(8)

    add_accent_box(slide, "💡 Sweet Spot: 2-4개 병렬", 7, left=4, width=8, color='warning')

def slide_11_warning(prs):
    """슬라이드 11: 비용 폭발 경고"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "⚠️ 비용 폭발 주의!", 0.8, font_size=44, color='error')

    warning = """실제 발생 사례 (AICosts.ai 보고서)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

• 49개 Sub-agents 병렬 실행
• 세션 시간: 2.5시간
• 토큰 소비: 887,000 tokens/분 (!!)
• 추정 비용: $8,000~$15,000 (단일 세션)

⚠️ 교훈: 병렬화는 강력하지만, 모니터링 없이는 위험!"""
    add_code_block(slide, warning, 2.2)

    add_accent_box(slide, "ccusage로 실시간 비용 모니터링 필수!", 6.5, left=2, width=12, color='warning')

def slide_12_prompt_examples(prs):
    """슬라이드 12: 실전 프롬프트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "실전 프롬프트 패턴", 0.8, font_size=44)

    prompt = """Use 4 parallel subagents to analyze this codebase:
- Agent 1: Analyze authentication in /lib/auth
- Agent 2: Analyze state management in /lib/providers
- Agent 3: Analyze API layer in /lib/services
- Agent 4: Analyze UI components in /lib/widgets

After all complete, synthesize into unified report."""
    add_code_block(slide, prompt, 2.2)

    add_accent_box(slide, "패턴: 분할 → 병렬 실행 → 종합", 6.5, left=3, width=10, color='primary')

def slide_13_optimization(prs):
    """슬라이드 13: 비용 최적화"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "비용 최적화 전략", 0.8, font_size=44)

    strategies = [
        "🎯 모델 믹싱: 탐색=Sonnet, 추론=Opus",
        "📁 파일 기반 통신: 결과를 .md에 저장 후 요약만 전달",
        "⏰ Quick Mode: 빠른 작업은 Haiku 활용"
    ]
    add_bullet_text(slide, strategies, 2.2, font_size=28)

    code = """# 파일 기반 통신 예시
"Each agent should write findings to /docs/tasks/{module}.md
Return only a 3-line summary. Do NOT pass full analysis."""
    add_code_block(slide, code, 5)

def slide_14_skills_concept(prs):
    """슬라이드 14: Skills 개념"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 4: Skills & Commands", 0.5, font_size=20, color='primary')
    add_title_text(slide, "개인 노하우 → 팀 자산으로", 1.2, font_size=44)

    transform = """개인의 노하우              →      팀 공용 Skills
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

"난 이렇게 해"             →      SKILL.md로 문서화
"물어보면 알려줄게"         →      자동으로 적용
"코드 리뷰 때 지적"         →      작성 시점에 반영
"퇴사하면 사라짐"           →      Git에 영구 보존"""
    add_code_block(slide, transform, 2.5)

def slide_15_skills_structure(prs):
    """슬라이드 15: Skills 폴더 구조"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Skills 폴더 구조", 0.8, font_size=44)

    structure = """.claude/
├── skills/
│   ├── flutter-conventions/
│   │   └── SKILL.md          # 코딩 컨벤션
│   ├── clean-architecture/
│   │   └── SKILL.md          # 아키텍처 패턴
│   └── testing-standards/
│       └── SKILL.md          # 테스트 작성 규칙
├── commands/
│   ├── feature.md            # 새 기능 워크플로우
│   ├── bugfix.md             # 버그 수정 워크플로우
│   └── review.md             # 코드 리뷰
└── CLAUDE.md                  # 프로젝트 컨텍스트"""
    add_code_block(slide, structure, 2)

    add_accent_box(slide, "💡 Git에 커밋하여 팀 전체 공유!", 7, left=3, width=10, color='success')

def slide_16_skill_example(prs):
    """슬라이드 16: SKILL.md 예시"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "SKILL.md 작성 예시", 0.8, font_size=44)

    example = """# Flutter Coding Conventions
triggers: ["flutter", "dart", "widget"]

## 네이밍 규칙
- 파일명: snake_case.dart
- 클래스명: PascalCase
- private: _underscorePrefix

## Widget 작성 규칙
// ✅ Good: const 생성자 사용
class MyWidget extends StatelessWidget {
  const MyWidget({super.key});
}

## 금지 사항
- print() 사용 금지 → logger.d() 사용
- dynamic 타입 사용 금지"""
    add_code_block(slide, example, 1.8)

def slide_17_commands(prs):
    """슬라이드 17: Custom Commands"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Custom Commands로 워크플로우 자동화", 0.8, font_size=44)

    commands = [
        "/feature → 새 기능 개발 A to Z 자동화",
        "/bugfix → 버그 조사 + 수정 + 테스트",
        "/review → 4가지 관점 병렬 코드 리뷰",
        "/release → 릴리즈 준비 자동화"
    ]

    for i, cmd in enumerate(commands):
        add_accent_box(slide, cmd, 2.2 + i * 1.4, left=1, width=14, color='primary')

    add_title_text(slide, "사용법: /feature implement user profile editing", 7.5, font_size=20, color='success')

def slide_18_mcp_concept(prs):
    """슬라이드 18: MCP 개념"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 5: MCP 통합", 0.5, font_size=20, color='primary')
    add_title_text(slide, "MCP = AI의 USB-C 포트", 1.2, font_size=44)

    diagram = """┌─────────────────────────────────────────────────┐
│  Claude Code ←──MCP──→ 외부 서비스               │
│      │                    │                      │
│      │                    ├─ Context7 (최신 문서) │
│      │                    ├─ Playwright (브라우저)│
│      │                    ├─ Notion (문서화)     │
│      │                    ├─ Linear (이슈 트래킹) │
│      │                    └─ GitHub (저장소)     │
└─────────────────────────────────────────────────┘

# 프로젝트 .mcp.json (Git 커밋 가능!)
→ 팀 전체가 동일한 MCP 환경 자동 적용"""
    add_code_block(slide, diagram, 2.2)

def slide_19_mcp_examples(prs):
    """슬라이드 19: MCP 활용 예시"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "4대 MCP 서버 활용법", 0.8, font_size=44)

    examples = [
        "Context7: 최신 Flutter/Riverpod API 문서 실시간 참조",
        "Playwright: E2E 테스트 자동 생성 및 실행",
        "Notion: 코드 리뷰 결과 자동 저장 및 문서 업데이트",
        "Linear: 버그 티켓 읽기 → 수정 → 상태 업데이트"
    ]

    for i, ex in enumerate(examples):
        color = ['primary', 'secondary', 'success', 'warning'][i]
        add_accent_box(slide, ex, 2.2 + i * 1.4, left=1, width=14, color=color)

def slide_20_mcp_workflow(prs):
    """슬라이드 20: MCP 통합 워크플로우"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "버그 수정 자동화 파이프라인", 0.8, font_size=44)

    workflow = """1. Linear MCP: 버그 티켓 PLA-123 읽기
        ↓
2. Context7 MCP: 관련 라이브러리 최신 정보 확인
        ↓
3. 코드 수정 + 테스트 작성 (Skills 적용)
        ↓
4. Playwright MCP: E2E 테스트로 수정 검증
        ↓
5. GitHub MCP: PR 생성
        ↓
6. Linear MCP: 티켓 상태 'In Review' 업데이트"""
    add_code_block(slide, workflow, 2)

    add_accent_box(slide, "터미널을 떠나지 않고 전체 워크플로우 자동화!", 7, left=2, width=12, color='success')

def slide_21_workflow(prs):
    """슬라이드 21: 새 기능 개발 워크플로우"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 6: 새 기능 개발 A to Z", 0.5, font_size=20, color='primary')
    add_title_text(slide, "8단계 자동화 프로세스", 1.2, font_size=44)

    steps = """1. /clear                    ← 깨끗한 시작
2. 3 explore agents          ← 코드베이스 분석
3. /ultrathink               ← 계획 수립
4. 3 parallel agents         ← Domain/Data/UI 동시 구현
5. testing-standards 적용     ← 테스트 생성
6. /review                   ← 품질 검증
7. notion + github MCP       ← 문서화 + PR
8. /handoff                  ← 다음 세션용 인수인계"""
    add_code_block(slide, steps, 2.2)

def slide_22_mindset(prs):
    """슬라이드 22: AI 시대 마인드셋"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 7: AI 시대 마인드셋", 0.5, font_size=20, color='primary')
    add_title_text(slide, "구현 → 설계에 시간 투자", 1.2, font_size=44)

    comparison = """Before AI:                    After AI:
├─ 구현: 80%                   ├─ 구현: 30%
├─ 설계: 15%                   ├─ 설계: 40%  ⬆️
└─ 리뷰: 5%                    ├─ 리뷰: 20%  ⬆️
                               └─ 취향/품질: 10%

집중해야 할 것:
✅ 좋은 시스템 설계
✅ 명확한 네이밍
✅ 철저한 문서화
✅ AI가 못하는 "왜?"에 대한 판단"""
    add_code_block(slide, comparison, 2.2)

def slide_23_levels(prs):
    """슬라이드 23: AI-Native 개발자 레벨"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "AI-Native 개발자 3단계", 0.8, font_size=44)

    levels = [
        ("Level 1: AI 사용자", "\"코드 작성해줘\" → 복사/붙여넣기", 'error'),
        ("Level 2: AI 협업자", "프롬프트 엔지니어링, 컨텍스트 관리", 'warning'),
        ("Level 3: AI 오케스트레이터 ← 목표!", "Skills/Commands + MCP + 병렬 처리 + 팀 지식 학습", 'success')
    ]

    for i, (title, desc, color) in enumerate(levels):
        add_accent_box(slide, title, 2 + i * 2, left=1, width=14, color=color)
        add_title_text(slide, desc, 3 + i * 2, font_size=20, color='text_secondary')

def slide_24_action(prs):
    """슬라이드 24: 액션 아이템"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "Part 8: 내일부터 바로 할 수 있는 것", 0.5, font_size=20, color='primary')
    add_title_text(slide, "Day 1-3 가이드", 1.2, font_size=44)

    actions = """Day 1: 기본 설정
/statusline show context percentage and model name
echo "# Project Context" > CLAUDE.md

Day 2: 첫 번째 커스텀 커맨드
mkdir -p .claude/commands
# review.md 템플릿 작성

Day 3: MCP 연결
claude mcp add context7 npx -y @context7/mcp-server
"use context7 to check latest Flutter features" """
    add_code_block(slide, actions, 2.2)

def slide_25_summary(prs):
    """슬라이드 25: 핵심 테이크어웨이 + Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    add_title_text(slide, "5가지 핵심 테이크어웨이", 0.8, font_size=44)

    takeaways = [
        "1️⃣ Context는 자산 → Status Line 모니터링, 50-60%에서 정리",
        "2️⃣ 병렬 처리 = 시간 구매 → 2-4개가 Sweet Spot",
        "3️⃣ Skills = 팀 지식 자산화 → Git에 영구 보존",
        "4️⃣ MCP = 생태계 통합 → 터미널에서 전체 워크플로우",
        "5️⃣ AI-Native 개발자 → 구현보다 설계와 취향에 투자"
    ]
    add_bullet_text(slide, takeaways, 2, font_size=24)

    add_title_text(slide, "Q&A", 7, font_size=36, color='primary')
    add_title_text(slide, "질문 있으시면 편하게 해주세요!", 7.7, font_size=20, color='text_secondary')

# ============================================
# MAIN
# ============================================

def main():
    prs = create_presentation()

    # Generate all slides
    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_why_hook(prs)
    slide_04_why_opus(prs)
    slide_05_context_concept(prs)
    slide_06_context_rot(prs)
    slide_07_statusline(prs)
    slide_08_context_action(prs)
    slide_09_subagent_concept(prs)
    slide_10_comparison(prs)
    slide_11_warning(prs)
    slide_12_prompt_examples(prs)
    slide_13_optimization(prs)
    slide_14_skills_concept(prs)
    slide_15_skills_structure(prs)
    slide_16_skill_example(prs)
    slide_17_commands(prs)
    slide_18_mcp_concept(prs)
    slide_19_mcp_examples(prs)
    slide_20_mcp_workflow(prs)
    slide_21_workflow(prs)
    slide_22_mindset(prs)
    slide_23_levels(prs)
    slide_24_action(prs)
    slide_25_summary(prs)

    # Save
    output_path = '/Users/woogi/Development/claude-craft/output/claude-code-seminar/Claude_Code_2.0_Seminar.pptx'
    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    print(f"📊 총 슬라이드: {len(prs.slides)}장")

if __name__ == "__main__":
    main()

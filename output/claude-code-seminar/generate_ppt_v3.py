#!/usr/bin/env python3
"""
Claude Code 2.0 세미나 PPT 생성 스크립트 v3.0
- 정밀한 레이아웃 시스템 적용
- Modern Dark 팔레트 (GitHub 스타일)
- 체계적인 그리드 시스템
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ============================================================
# 설계 상수 (SPEC.md 기반)
# ============================================================

# 슬라이드 크기
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# 여백
MARGIN_TOP = 0.5
MARGIN_BOTTOM = 0.5
MARGIN_LEFT = 0.7
MARGIN_RIGHT = 0.7

# 콘텐츠 영역
CONTENT_X = MARGIN_LEFT
CONTENT_Y = 1.2  # 헤더 아래
CONTENT_W = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT  # 11.933
CONTENT_H = 5.1  # 콘텐츠 영역 높이

# 헤더/푸터
HEADER_Y = 0.5
FOOTER_Y = 6.5

# ============================================================
# 컬러 팔레트: Modern Dark (GitHub 스타일)
# ============================================================

class Colors:
    # 배경
    BG_PRIMARY = RGBColor(0x0d, 0x11, 0x17)     # #0d1117
    BG_SECONDARY = RGBColor(0x16, 0x1b, 0x22)   # #161b22
    BG_CARD = RGBColor(0x21, 0x26, 0x2d)        # #21262d

    # 텍스트
    TEXT_PRIMARY = RGBColor(0xf0, 0xf6, 0xfc)   # #f0f6fc
    TEXT_SECONDARY = RGBColor(0x8b, 0x94, 0x9e) # #8b949e
    TEXT_MUTED = RGBColor(0x6e, 0x76, 0x81)     # #6e7681

    # 강조
    ACCENT_BLUE = RGBColor(0x58, 0xa6, 0xff)    # #58a6ff
    ACCENT_GREEN = RGBColor(0x3f, 0xb9, 0x50)   # #3fb950
    ACCENT_PURPLE = RGBColor(0xa3, 0x71, 0xf7)  # #a371f7
    ACCENT_YELLOW = RGBColor(0xd2, 0x99, 0x22)  # #d29922
    ACCENT_RED = RGBColor(0xf8, 0x51, 0x49)     # #f85149

# ============================================================
# 타이포그래피
# ============================================================

class FontSize:
    HERO = 44
    SECTION = 36
    TITLE = 28
    SUBTITLE = 20
    BODY = 16
    CAPTION = 12
    LABEL = 10

# ============================================================
# 유틸리티 함수
# ============================================================

def create_presentation():
    """16:9 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    return prs

def add_slide(prs):
    """빈 슬라이드 추가 (배경색 포함)"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.BG_PRIMARY

    return slide

def add_text(slide, text, left, top, width, height,
             font_size=FontSize.BODY, bold=False, italic=False,
             color=None, align='left', valign='top', font_name='Arial'):
    """텍스트 박스 추가"""
    if color is None:
        color = Colors.TEXT_PRIMARY

    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    # 수직 정렬
    if valign == 'middle':
        tf.anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name

    # 수평 정렬
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    return text_box

def add_shape(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    """도형 추가"""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 테두리 제거
    return shape

def add_bullets(slide, items, left, top, width, height,
                font_size=FontSize.BODY, color=None):
    """불릿 리스트 추가"""
    if color is None:
        color = Colors.TEXT_SECONDARY

    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"•  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(14)
        p.level = 0

    return text_box

def add_badge(slide, text, left=MARGIN_LEFT, top=HEADER_Y):
    """섹션 뱃지 추가"""
    badge_w = 1.6
    badge_h = 0.35

    # 배경
    add_shape(slide, left, top, badge_w, badge_h, Colors.ACCENT_BLUE)

    # 텍스트
    add_text(slide, text.upper(),
             left, top, badge_w, badge_h,
             font_size=FontSize.LABEL, bold=True,
             color=Colors.TEXT_PRIMARY, align='center', valign='middle')

def add_page_number(slide, current, total=25):
    """페이지 번호 추가"""
    page_w = 1.2
    add_text(slide, f"{current:02d} / {total:02d}",
             SLIDE_WIDTH - MARGIN_RIGHT - page_w, HEADER_Y,
             page_w, 0.35,
             font_size=FontSize.CAPTION,
             color=Colors.TEXT_MUTED, align='right', valign='middle')

# ============================================================
# 슬라이드 템플릿
# ============================================================

def slide_cover(prs, title, subtitle, presenter="", date=""):
    """Cover Slide (표지)"""
    slide = add_slide(prs)

    # 상단 악센트 바
    add_shape(slide, 0, 0, SLIDE_WIDTH, 0.15, Colors.ACCENT_BLUE)

    # 타이틀
    add_text(slide, title,
             MARGIN_LEFT, 2.6, CONTENT_W, 1.2,
             font_size=FontSize.HERO, bold=True,
             color=Colors.TEXT_PRIMARY, align='center', valign='middle')

    # 구분선
    add_shape(slide, 4.0, 3.9, 5.333, 0.02, Colors.ACCENT_BLUE)

    # 부제목
    add_text(slide, subtitle,
             MARGIN_LEFT, 4.1, CONTENT_W, 0.6,
             font_size=FontSize.SUBTITLE,
             color=Colors.TEXT_SECONDARY, align='center')

    # 발표자 (좌하단)
    if presenter:
        add_text(slide, presenter,
                 MARGIN_LEFT, FOOTER_Y, 4.0, 0.4,
                 font_size=FontSize.CAPTION,
                 color=Colors.TEXT_MUTED, align='left')

    # 날짜 (우하단)
    if date:
        add_text(slide, date,
                 SLIDE_WIDTH - MARGIN_RIGHT - 4.0, FOOTER_Y, 4.0, 0.4,
                 font_size=FontSize.CAPTION,
                 color=Colors.TEXT_MUTED, align='right')

    return slide

def slide_contents(prs, sections):
    """Contents Slide (목차)"""
    slide = add_slide(prs)

    # 타이틀
    add_text(slide, "CONTENTS",
             MARGIN_LEFT, MARGIN_TOP, CONTENT_W, 0.8,
             font_size=FontSize.SECTION, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 구분선
    add_shape(slide, MARGIN_LEFT, 1.3, 2.0, 0.03, Colors.ACCENT_BLUE)

    # 섹션 리스트
    start_y = 1.8
    line_height = 0.6

    for i, section in enumerate(sections):
        y = start_y + (i * line_height)

        # 번호
        add_text(slide, f"{i+1:02d}",
                 MARGIN_LEFT, y, 0.6, 0.5,
                 font_size=FontSize.SUBTITLE, bold=True,
                 color=Colors.ACCENT_BLUE, valign='middle')

        # 제목
        add_text(slide, section['title'],
                 MARGIN_LEFT + 0.8, y, 8.0, 0.5,
                 font_size=FontSize.BODY,
                 color=Colors.TEXT_PRIMARY, valign='middle')

        # 페이지 번호
        add_text(slide, f"{section['page']:02d}",
                 SLIDE_WIDTH - MARGIN_RIGHT - 0.6, y, 0.6, 0.5,
                 font_size=FontSize.BODY,
                 color=Colors.TEXT_MUTED, align='right', valign='middle')

    return slide

def slide_section(prs, number, title, description=""):
    """Section Divider (섹션 구분)"""
    slide = add_slide(prs)

    # 배경을 악센트 블루로
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.ACCENT_BLUE

    # 섹션 번호
    add_text(slide, f"{number:02d}",
             MARGIN_LEFT, 2.0, 2.0, 1.2,
             font_size=72, bold=True,
             color=Colors.TEXT_PRIMARY, valign='middle')

    # 섹션 타이틀
    add_text(slide, title,
             MARGIN_LEFT, 3.4, CONTENT_W, 0.9,
             font_size=FontSize.SECTION, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 설명
    if description:
        add_text(slide, description,
                 MARGIN_LEFT, 4.5, CONTENT_W, 0.6,
                 font_size=FontSize.BODY,
                 color=Colors.BG_PRIMARY)  # 어두운 텍스트

    return slide

def slide_content(prs, headline, bullets, section="", page=0):
    """Content Slide (기본 콘텐츠)"""
    slide = add_slide(prs)

    if section:
        add_badge(slide, section)
    if page:
        add_page_number(slide, page)

    # 헤드라인
    add_text(slide, headline,
             MARGIN_LEFT, CONTENT_Y, CONTENT_W, 0.7,
             font_size=FontSize.TITLE, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 불릿 포인트
    add_bullets(slide, bullets,
                MARGIN_LEFT, CONTENT_Y + 0.9, CONTENT_W, 4.0)

    return slide

def slide_two_column(prs, headline, bullets, section="", page=0):
    """Two Column Slide (2단 레이아웃)"""
    slide = add_slide(prs)

    if section:
        add_badge(slide, section)
    if page:
        add_page_number(slide, page)

    # 헤드라인
    add_text(slide, headline,
             MARGIN_LEFT, CONTENT_Y, CONTENT_W, 0.7,
             font_size=FontSize.TITLE, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 왼쪽 컬럼: 불릿
    col_width = 5.2
    add_bullets(slide, bullets,
                MARGIN_LEFT, CONTENT_Y + 0.9, col_width, 4.0)

    # 오른쪽 컬럼: 비주얼 플레이스홀더
    right_x = MARGIN_LEFT + col_width + 0.5
    right_w = CONTENT_W - col_width - 0.5

    add_shape(slide, right_x, CONTENT_Y + 0.9, right_w, 3.5, Colors.BG_CARD)
    add_text(slide, "[ Diagram / Visual ]",
             right_x, CONTENT_Y + 2.4, right_w, 0.5,
             font_size=FontSize.CAPTION,
             color=Colors.TEXT_MUTED, align='center', valign='middle')

    return slide

def slide_stats(prs, headline, metrics, section="", page=0):
    """Statistics Slide (통계)"""
    slide = add_slide(prs)

    if section:
        add_badge(slide, section)
    if page:
        add_page_number(slide, page)

    # 헤드라인
    add_text(slide, headline,
             MARGIN_LEFT, CONTENT_Y, CONTENT_W, 0.7,
             font_size=FontSize.TITLE, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 카드 레이아웃
    num_cards = len(metrics)
    card_w = 3.4
    card_h = 2.8
    gap = 0.5
    total_w = (num_cards * card_w) + ((num_cards - 1) * gap)
    start_x = (SLIDE_WIDTH - total_w) / 2
    card_y = 2.3

    for i, metric in enumerate(metrics):
        x = start_x + (i * (card_w + gap))

        # 카드 배경
        add_shape(slide, x, card_y, card_w, card_h, Colors.BG_CARD)

        # 값
        add_text(slide, metric['value'],
                 x, card_y + 0.5, card_w, 1.2,
                 font_size=36, bold=True,
                 color=Colors.ACCENT_GREEN, align='center', valign='middle')

        # 라벨
        add_text(slide, metric['label'],
                 x, card_y + 2.0, card_w, 0.5,
                 font_size=FontSize.BODY,
                 color=Colors.TEXT_SECONDARY, align='center', valign='middle')

    return slide

def slide_comparison(prs, headline, before_items, after_items, section="", page=0):
    """Comparison Slide (비교)"""
    slide = add_slide(prs)

    if section:
        add_badge(slide, section)
    if page:
        add_page_number(slide, page)

    # 헤드라인
    add_text(slide, headline,
             MARGIN_LEFT, CONTENT_Y, CONTENT_W, 0.7,
             font_size=FontSize.TITLE, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 컬럼 설정
    col_w = 5.6
    gap = 0.5
    left_x = MARGIN_LEFT
    right_x = MARGIN_LEFT + col_w + gap
    box_y = CONTENT_Y + 0.9
    box_h = 3.8
    header_h = 0.45

    # BEFORE 컬럼
    add_shape(slide, left_x, box_y, col_w, header_h, Colors.ACCENT_RED)
    add_text(slide, "BEFORE",
             left_x, box_y, col_w, header_h,
             font_size=FontSize.LABEL, bold=True,
             color=Colors.TEXT_PRIMARY, align='center', valign='middle')
    add_bullets(slide, before_items,
                left_x + 0.2, box_y + header_h + 0.3, col_w - 0.4, box_h - header_h - 0.5)

    # AFTER 컬럼
    add_shape(slide, right_x, box_y, col_w, header_h, Colors.ACCENT_GREEN)
    add_text(slide, "AFTER",
             right_x, box_y, col_w, header_h,
             font_size=FontSize.LABEL, bold=True,
             color=Colors.BG_PRIMARY, align='center', valign='middle')
    add_bullets(slide, after_items,
                right_x + 0.2, box_y + header_h + 0.3, col_w - 0.4, box_h - header_h - 0.5)

    return slide

def slide_code(prs, headline, code_lines, section="", page=0):
    """Code Slide (코드 블록)"""
    slide = add_slide(prs)

    if section:
        add_badge(slide, section)
    if page:
        add_page_number(slide, page)

    # 헤드라인
    add_text(slide, headline,
             MARGIN_LEFT, CONTENT_Y, CONTENT_W, 0.7,
             font_size=FontSize.TITLE, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 코드 블록 배경
    code_y = CONTENT_Y + 0.9
    code_h = 3.8
    add_shape(slide, MARGIN_LEFT, code_y, CONTENT_W, code_h, Colors.BG_CARD)

    # 코드 텍스트
    text_box = slide.shapes.add_textbox(
        Inches(MARGIN_LEFT + 0.3), Inches(code_y + 0.3),
        Inches(CONTENT_W - 0.6), Inches(code_h - 0.6)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = Pt(14)
        p.font.name = 'Consolas'
        p.font.color.rgb = Colors.ACCENT_GREEN
        p.space_after = Pt(6)

    return slide

def slide_action(prs, headline, actions, section="", page=0):
    """Action Items Slide"""
    slide = add_slide(prs)

    if section:
        add_badge(slide, section)
    if page:
        add_page_number(slide, page)

    # 헤드라인
    add_text(slide, headline,
             MARGIN_LEFT, CONTENT_Y, CONTENT_W, 0.7,
             font_size=FontSize.TITLE, bold=True,
             color=Colors.TEXT_PRIMARY)

    # 액션 카드
    card_h = 0.8
    gap = 0.2
    start_y = CONTENT_Y + 1.0

    for i, action in enumerate(actions):
        y = start_y + (i * (card_h + gap))

        # 카드 배경
        add_shape(slide, MARGIN_LEFT, y, CONTENT_W, card_h, Colors.BG_CARD)

        # 번호
        num_w = 0.5
        add_shape(slide, MARGIN_LEFT + 0.15, y + 0.15, num_w, num_w, Colors.ACCENT_BLUE)
        add_text(slide, str(i + 1),
                 MARGIN_LEFT + 0.15, y + 0.15, num_w, num_w,
                 font_size=FontSize.BODY, bold=True,
                 color=Colors.TEXT_PRIMARY, align='center', valign='middle')

        # 액션 텍스트
        add_text(slide, action,
                 MARGIN_LEFT + 0.9, y, CONTENT_W - 1.1, card_h,
                 font_size=FontSize.BODY,
                 color=Colors.TEXT_PRIMARY, valign='middle')

    return slide

def slide_closing(prs, title, subtitle="", contact=""):
    """Closing Slide (마무리)"""
    slide = add_slide(prs)

    # 배경을 악센트 블루로
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.ACCENT_BLUE

    # 타이틀
    add_text(slide, title,
             MARGIN_LEFT, 2.5, CONTENT_W, 1.2,
             font_size=FontSize.HERO, bold=True,
             color=Colors.TEXT_PRIMARY, align='center', valign='middle')

    # 부제목
    if subtitle:
        add_text(slide, subtitle,
                 MARGIN_LEFT, 3.9, CONTENT_W, 0.6,
                 font_size=FontSize.SUBTITLE,
                 color=Colors.BG_PRIMARY, align='center')

    # 연락처
    if contact:
        add_text(slide, contact,
                 MARGIN_LEFT, 5.2, CONTENT_W, 0.8,
                 font_size=FontSize.BODY,
                 color=Colors.BG_PRIMARY, align='center')

    return slide

# ============================================================
# 메인: Claude Code 2.0 세미나 PPT 생성
# ============================================================

def generate_seminar():
    """Claude Code 2.0 세미나 PPT 생성"""
    prs = create_presentation()

    # ──────────────────────────────────────────────────────
    # 1. Cover
    # ──────────────────────────────────────────────────────
    slide_cover(prs,
        title="Claude Code 2.0\n완벽 활용 가이드",
        subtitle="AI 시대, 개발 생산성을 극대화하는 방법",
        presenter="Tech Team",
        date="2025.01"
    )

    # ──────────────────────────────────────────────────────
    # 2. Contents
    # ──────────────────────────────────────────────────────
    slide_contents(prs, [
        {'title': 'Why Claude Code', 'page': 3},
        {'title': 'Context Engineering', 'page': 5},
        {'title': '병렬 Sub-agents', 'page': 9},
        {'title': 'Skills & Commands', 'page': 14},
        {'title': 'MCP 통합', 'page': 18},
        {'title': '실전 & 마인드셋', 'page': 21},
        {'title': 'Closing', 'page': 24},
    ])

    # ──────────────────────────────────────────────────────
    # Part 1: Why Claude Code (3-4)
    # ──────────────────────────────────────────────────────
    slide_section(prs, 1, "Why Claude Code",
                  "AI 코딩 어시스턴트의 새로운 기준")

    slide_stats(prs,
        headline="개발 속도 2-3배 향상의 비밀",
        metrics=[
            {'value': '2-3x', 'label': '개발 속도'},
            {'value': '200K', 'label': 'Context Window'},
            {'value': 'Opus 4.5', 'label': '최신 모델'},
        ],
        section="WHY", page=4
    )

    # ──────────────────────────────────────────────────────
    # Part 2: Context Engineering (5-8)
    # ──────────────────────────────────────────────────────
    slide_section(prs, 2, "Context Engineering",
                  "200K 토큰 공간을 효과적으로 관리하는 방법")

    slide_two_column(prs,
        headline="Context Window = 200K 토큰 공간",
        bullets=[
            "시스템 프롬프트 + 대화 이력 + 파일 내용",
            "토큰 ≈ 단어의 3/4 (영어 기준)",
            "한글은 토큰 소모가 더 큼",
            "200K ≈ 소설 2-3권 분량",
        ],
        section="CONTEXT", page=6
    )

    slide_content(prs,
        headline="Context Rot - 성능 저하의 원인",
        bullets=[
            "Context가 80% 이상 차면 성능 급격히 저하",
            "오래된 정보가 최신 정보와 혼란 유발",
            "불필요한 파일 읽기로 공간 낭비",
            "해결책: /compact 명령어로 정리",
        ],
        section="CONTEXT", page=7
    )

    slide_code(prs,
        headline="Status Line으로 실시간 모니터링",
        code_lines=[
            '# .claude/settings.json',
            '{',
            '  "statusLineEnabled": true,',
            '  "statusLineTemplate":',
            '    "[{model}] {tokens}/{max} ({percent}%)"',
            '}'
        ],
        section="CONTEXT", page=8
    )

    # ──────────────────────────────────────────────────────
    # Part 3: 병렬 Sub-agents (9-13)
    # ──────────────────────────────────────────────────────
    slide_section(prs, 3, "병렬 Sub-agents",
                  "독립 Context의 분신들로 시간 단축")

    slide_two_column(prs,
        headline="Sub-agent = 독립 Context의 분신들",
        bullets=[
            "메인 Agent와 별도의 Context 유지",
            "병렬 실행으로 작업 시간 단축",
            "각자 다른 파일/작업 동시 처리",
            "결과만 메인 Agent로 전달",
        ],
        section="AGENTS", page=10
    )

    slide_comparison(prs,
        headline="순차 vs 병렬: 시간 62% 단축",
        before_items=[
            "파일 1 분석 → 완료",
            "파일 2 분석 → 완료",
            "파일 3 분석 → 완료",
            "총 소요시간: 30분",
        ],
        after_items=[
            "파일 1, 2, 3 동시 분석",
            "결과 통합",
            "총 소요시간: 12분",
            "62% 시간 절약!",
        ],
        section="AGENTS", page=11
    )

    slide_content(prs,
        headline="비용 폭발 주의! 49개 병렬의 교훈",
        bullets=[
            "병렬 Agent 수 = 비용 배수",
            "권장: 2-4개가 Sweet Spot",
            "5개 이상은 수익 체감",
            "모델 믹싱으로 비용 최적화 (Haiku 활용)",
        ],
        section="AGENTS", page=12
    )

    slide_code(prs,
        headline="실전 프롬프트 패턴",
        code_lines=[
            '# 병렬 분석 요청',
            '"다음 3개 파일을 병렬로 분석해줘:',
            '  - src/auth.ts',
            '  - src/api.ts',
            '  - src/db.ts',
            '각각 보안 취약점을 찾아줘"',
        ],
        section="AGENTS", page=13
    )

    # ──────────────────────────────────────────────────────
    # Part 4: Skills & Commands (14-17)
    # ──────────────────────────────────────────────────────
    slide_section(prs, 4, "Skills & Commands",
                  "개인 노하우를 팀 자산으로 변환")

    slide_two_column(prs,
        headline="Skills = 재사용 가능한 지식 모듈",
        bullets=[
            "SKILL.md 파일에 워크플로우 정의",
            ".claude/skills/ 폴더에 저장",
            "Git으로 팀 전체 공유 가능",
            "복잡한 작업을 명령어 하나로!",
        ],
        section="SKILLS", page=15
    )

    slide_code(prs,
        headline="SKILL.md 작성 예시",
        code_lines=[
            '---',
            'name: code-review',
            'description: 코드 리뷰 자동화',
            '---',
            '',
            '# Code Review Skill',
            '1. 파일 변경사항 분석',
            '2. 코딩 컨벤션 체크',
            '3. 보안 취약점 스캔',
        ],
        section="SKILLS", page=16
    )

    slide_content(prs,
        headline="Custom Commands로 워크플로우 자동화",
        bullets=[
            "/feature - 새 기능 개발 시작",
            "/review - 코드 리뷰 실행",
            "/test - 테스트 작성 및 실행",
            "/deploy - 배포 프로세스 시작",
        ],
        section="SKILLS", page=17
    )

    # ──────────────────────────────────────────────────────
    # Part 5: MCP 통합 (18-20)
    # ──────────────────────────────────────────────────────
    slide_section(prs, 5, "MCP 통합",
                  "AI의 USB-C 포트 - 무한 확장 가능")

    slide_two_column(prs,
        headline="MCP = Model Context Protocol",
        bullets=[
            "AI와 외부 시스템 연결 표준",
            "데이터베이스, API, 파일시스템",
            "플러그인처럼 기능 확장",
            "한 번 설정, 계속 사용",
        ],
        section="MCP", page=19
    )

    slide_stats(prs,
        headline="4대 MCP 서버 활용법",
        metrics=[
            {'value': 'Context7', 'label': '문서 검색'},
            {'value': 'Playwright', 'label': '브라우저'},
            {'value': 'SQLite', 'label': '데이터베이스'},
        ],
        section="MCP", page=20
    )

    # ──────────────────────────────────────────────────────
    # Part 6: 실전 & 마인드셋 (21-23)
    # ──────────────────────────────────────────────────────
    slide_section(prs, 6, "실전 & 마인드셋",
                  "AI-Native 개발자로 성장하기")

    slide_content(prs,
        headline="새 기능 개발 A to Z",
        bullets=[
            "1. 요구사항 정리 (프롬프트 작성)",
            "2. 설계 리뷰 (AI와 토론)",
            "3. 코드 생성 (병렬 작업)",
            "4. 테스트 & 리팩토링",
            "5. 코드 리뷰 & 배포",
        ],
        section="MINDSET", page=22
    )

    slide_comparison(prs,
        headline="구현 → 설계에 시간 투자",
        before_items=[
            "구현 70%",
            "설계 15%",
            "테스트 15%",
            "야근이 일상...",
        ],
        after_items=[
            "설계 50%",
            "구현 20%",
            "검증 30%",
            "워라밸 확보!",
        ],
        section="MINDSET", page=23
    )

    # ──────────────────────────────────────────────────────
    # Closing (24-25)
    # ──────────────────────────────────────────────────────
    slide_action(prs,
        headline="내일부터 바로 할 수 있는 3가지",
        actions=[
            "Day 1: Status Line 설정하고 Context 모니터링 시작",
            "Day 2: 반복 작업을 Skill로 만들어보기",
            "Day 3: 병렬 Agent로 코드 리뷰 자동화",
        ],
        section="ACTION", page=24
    )

    slide_closing(prs,
        title="Thank You",
        subtitle="Questions & Discussion",
        contact="AI가 코드를 짜는 시대, 우리는 더 나은 설계자가 됩니다."
    )

    # ──────────────────────────────────────────────────────
    # 파일 저장
    # ──────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__),
                               "Claude_Code_2.0_Seminar_v3.pptx")
    prs.save(output_path)

    print(f"✅ PPT 생성 완료: {output_path}")
    print(f"📊 총 슬라이드: {len(prs.slides)}장")

    return output_path

if __name__ == "__main__":
    generate_seminar()

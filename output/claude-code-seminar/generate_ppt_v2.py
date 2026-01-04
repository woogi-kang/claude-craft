#!/usr/bin/env python3
"""
Claude Code 2.0 세미나 PPT 생성 스크립트 v2.0
- 정밀한 타이포그래피 시스템 적용
- Modern Dark 팔레트 사용
- 10가지 슬라이드 템플릿 구현
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

# ============================================================
# 디자인 시스템: Modern Dark 팔레트
# ============================================================

PALETTE = {
    'bg_primary': '0f0f0f',       # 메인 배경
    'bg_secondary': '1a1a1a',     # 보조 배경
    'bg_card': '252525',          # 카드 배경
    'text_primary': 'ffffff',     # 제목
    'text_secondary': 'a0a0a0',   # 본문
    'accent': '667eea',           # 강조색 (파랑)
    'accent_secondary': '764ba2', # 보조 강조색 (보라)
    'success': '4ecdc4',          # 성공
    'warning': 'ffe66d',          # 경고
    'error': 'ff6b6b',            # 에러
}

# ============================================================
# 타이포그래피 시스템
# ============================================================

TYPOGRAPHY = {
    'hero': {'size': 54, 'bold': True},        # 타이틀 슬라이드 (84pt → 54pt 조정)
    'section': {'size': 44, 'bold': True},     # 섹션 구분
    'title': {'size': 32, 'bold': True},       # 슬라이드 제목
    'subtitle': {'size': 20, 'bold': False},   # 부제목
    'body': {'size': 16, 'bold': False},       # 본문
    'caption': {'size': 12, 'bold': False},    # 캡션
    'label': {'size': 10, 'bold': True},       # 라벨
    'metric': {'size': 48, 'bold': True},      # 통계 수치
}

# ============================================================
# 레이아웃 상수 (인치 단위)
# ============================================================

SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 5.625
PADDING = 0.67          # 48pt / 72
CONTENT_WIDTH = SLIDE_WIDTH - (PADDING * 2)

# ============================================================
# 유틸리티 함수
# ============================================================

def hex_to_rgb(hex_color):
    """HEX 색상을 RGBColor로 변환"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def create_presentation():
    """16:9 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    return prs

def add_slide(prs):
    """빈 슬라이드 추가"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_background(slide, color_key='bg_primary'):
    """배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(PALETTE[color_key])

def add_shape_with_fill(slide, left, top, width, height, color_key, radius=0):
    """색상이 채워진 도형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(PALETTE[color_key])
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 style='body', color_key='text_primary', align='left',
                 valign='top'):
    """텍스트 박스 추가 (타이포그래피 시스템 적용)"""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    # 수직 정렬
    if valign == 'middle':
        tf.anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text

    # 타이포그래피 적용
    typo = TYPOGRAPHY.get(style, TYPOGRAPHY['body'])
    p.font.size = Pt(typo['size'])
    p.font.bold = typo['bold']
    p.font.color.rgb = hex_to_rgb(PALETTE[color_key])
    p.font.name = 'Arial'  # 시스템 폰트 폴백

    # 정렬
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    return text_box

def add_bullet_list(slide, items, left, top, width, height,
                    color_key='text_secondary', font_size=16):
    """불릿 리스트 추가"""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_to_rgb(PALETTE[color_key])
        p.font.name = 'Arial'
        p.space_after = Pt(12)

    return text_box

def add_section_badge(slide, text, left=PADDING, top=0.4):
    """섹션 뱃지 추가"""
    badge_width = len(text) * 0.12 + 0.4

    # 배경 도형
    shape = add_shape_with_fill(slide, left, top, badge_width, 0.35, 'accent')

    # 텍스트
    add_text_box(slide, text, left, top, badge_width, 0.35,
                 style='label', color_key='text_primary', align='center', valign='middle')

def add_page_number(slide, num, total=25):
    """페이지 번호 추가"""
    add_text_box(slide, f"{num:02d}/{total:02d}",
                 left=SLIDE_WIDTH - PADDING - 0.8, top=0.4,
                 width=0.8, height=0.35,
                 style='caption', color_key='text_secondary', align='right')

# ============================================================
# 슬라이드 템플릿
# ============================================================

def create_cover_slide(prs, title, subtitle, presenter="", date=""):
    """1. Cover Slide (표지)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    # 상단 그라데이션 효과 (단색 바로 대체)
    add_shape_with_fill(slide, 0, 0, SLIDE_WIDTH, 0.3, 'accent')

    # Hero Title
    add_text_box(slide, title,
                 left=PADDING, top=1.8, width=CONTENT_WIDTH, height=1.2,
                 style='hero', color_key='text_primary', align='center', valign='middle')

    # Subtitle
    add_text_box(slide, subtitle,
                 left=PADDING, top=3.1, width=CONTENT_WIDTH, height=0.5,
                 style='subtitle', color_key='text_secondary', align='center')

    # Footer
    if presenter or date:
        footer_text = f"{presenter}  |  {date}" if presenter and date else presenter or date
        add_text_box(slide, footer_text,
                     left=PADDING, top=4.8, width=CONTENT_WIDTH, height=0.3,
                     style='caption', color_key='text_secondary', align='center')

    return slide

def create_contents_slide(prs, sections):
    """2. Contents Slide (목차)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    # Title
    add_text_box(slide, "CONTENTS",
                 left=PADDING, top=PADDING, width=CONTENT_WIDTH, height=0.6,
                 style='title', color_key='text_primary')

    # Section List
    start_y = 1.4
    for i, section in enumerate(sections):
        y = start_y + (i * 0.55)

        # Section Number
        add_text_box(slide, f"{i+1:02d}",
                     left=PADDING, top=y, width=0.5, height=0.4,
                     style='subtitle', color_key='accent', align='left')

        # Section Title
        add_text_box(slide, section['title'],
                     left=PADDING + 0.7, top=y, width=6.0, height=0.4,
                     style='body', color_key='text_primary')

        # Page Number
        add_text_box(slide, f"{section['page']:02d}",
                     left=SLIDE_WIDTH - PADDING - 0.5, top=y, width=0.5, height=0.4,
                     style='body', color_key='text_secondary', align='right')

    return slide

def create_section_divider(prs, section_num, title, description=""):
    """3. Section Divider (섹션 구분)"""
    slide = add_slide(prs)
    set_background(slide, 'accent')

    # Section Number
    add_text_box(slide, f"{section_num:02d}",
                 left=PADDING, top=1.5, width=1.5, height=0.8,
                 style='section', color_key='text_primary')

    # Title
    add_text_box(slide, title,
                 left=PADDING, top=2.3, width=CONTENT_WIDTH, height=0.8,
                 style='section', color_key='text_primary')

    # Description
    if description:
        add_text_box(slide, description,
                     left=PADDING, top=3.2, width=CONTENT_WIDTH, height=0.5,
                     style='body', color_key='text_primary')

    return slide

def create_content_slide(prs, headline, bullets, section="", page_num=0,
                         visual_placeholder=False):
    """4. Content Slide (콘텐츠)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    # Section Badge
    if section:
        add_section_badge(slide, section)

    # Page Number
    if page_num:
        add_page_number(slide, page_num)

    # Headline
    add_text_box(slide, headline,
                 left=PADDING, top=1.0, width=CONTENT_WIDTH, height=0.6,
                 style='title', color_key='text_primary')

    # Content Area
    if visual_placeholder:
        # 2-column layout
        content_width = 4.0
        add_bullet_list(slide, bullets,
                        left=PADDING, top=1.8, width=content_width, height=3.2)

        # Visual Placeholder
        add_shape_with_fill(slide, 5.2, 1.8, 4.1, 3.0, 'bg_card')
        add_text_box(slide, "[Visual]",
                     left=5.2, top=3.0, width=4.1, height=0.5,
                     style='caption', color_key='text_secondary', align='center')
    else:
        # Full width bullets
        add_bullet_list(slide, bullets,
                        left=PADDING, top=1.8, width=CONTENT_WIDTH, height=3.2)

    return slide

def create_statistics_slide(prs, title, metrics, section="", page_num=0):
    """5. Statistics Slide (통계/수치)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    # Section Badge
    if section:
        add_section_badge(slide, section)

    # Page Number
    if page_num:
        add_page_number(slide, page_num)

    # Title
    add_text_box(slide, title,
                 left=PADDING, top=1.0, width=CONTENT_WIDTH, height=0.6,
                 style='title', color_key='text_primary')

    # Metric Cards
    num_metrics = len(metrics)
    card_width = 2.5
    gap = 0.4
    total_width = (num_metrics * card_width) + ((num_metrics - 1) * gap)
    start_x = (SLIDE_WIDTH - total_width) / 2

    for i, metric in enumerate(metrics):
        x = start_x + (i * (card_width + gap))

        # Card Background
        add_shape_with_fill(slide, x, 1.9, card_width, 2.4, 'bg_card')

        # Metric Value
        add_text_box(slide, metric['value'],
                     left=x, top=2.2, width=card_width, height=1.0,
                     style='metric', color_key='accent', align='center', valign='middle')

        # Metric Label
        add_text_box(slide, metric['label'],
                     left=x, top=3.5, width=card_width, height=0.5,
                     style='body', color_key='text_secondary', align='center')

    return slide

def create_comparison_slide(prs, title, before_items, after_items,
                           section="", page_num=0):
    """6. Comparison Slide (비교)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    if section:
        add_section_badge(slide, section)
    if page_num:
        add_page_number(slide, page_num)

    # Title
    add_text_box(slide, title,
                 left=PADDING, top=1.0, width=CONTENT_WIDTH, height=0.6,
                 style='title', color_key='text_primary')

    col_width = 4.0
    left_x = PADDING
    right_x = SLIDE_WIDTH / 2 + 0.2

    # Before Column
    add_shape_with_fill(slide, left_x, 1.8, col_width, 0.4, 'error')
    add_text_box(slide, "BEFORE",
                 left=left_x, top=1.8, width=col_width, height=0.4,
                 style='label', color_key='text_primary', align='center', valign='middle')
    add_bullet_list(slide, before_items,
                    left=left_x, top=2.4, width=col_width, height=2.5)

    # After Column
    add_shape_with_fill(slide, right_x, 1.8, col_width, 0.4, 'success')
    add_text_box(slide, "AFTER",
                 left=right_x, top=1.8, width=col_width, height=0.4,
                 style='label', color_key='bg_primary', align='center', valign='middle')
    add_bullet_list(slide, after_items,
                    left=right_x, top=2.4, width=col_width, height=2.5)

    return slide

def create_quote_slide(prs, quote, attribution="", role=""):
    """7. Quote Slide (인용)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_secondary')

    # Quote Mark
    add_text_box(slide, '"',
                 left=PADDING, top=1.2, width=1.0, height=0.8,
                 style='hero', color_key='accent')

    # Quote Text
    add_text_box(slide, quote,
                 left=PADDING + 0.3, top=1.8, width=CONTENT_WIDTH - 0.6, height=2.0,
                 style='subtitle', color_key='text_primary', align='center', valign='middle')

    # Attribution
    if attribution:
        attr_text = f"— {attribution}"
        if role:
            attr_text += f"\n{role}"
        add_text_box(slide, attr_text,
                     left=PADDING, top=4.0, width=CONTENT_WIDTH, height=0.8,
                     style='body', color_key='text_secondary', align='center')

    return slide

def create_code_slide(prs, title, code_lines, section="", page_num=0):
    """8. Code Slide (코드 예시)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    if section:
        add_section_badge(slide, section)
    if page_num:
        add_page_number(slide, page_num)

    # Title
    add_text_box(slide, title,
                 left=PADDING, top=1.0, width=CONTENT_WIDTH, height=0.6,
                 style='title', color_key='text_primary')

    # Code Block Background
    add_shape_with_fill(slide, PADDING, 1.7, CONTENT_WIDTH, 3.2, 'bg_card')

    # Code Text
    code_text = '\n'.join(code_lines)
    text_box = slide.shapes.add_textbox(
        Inches(PADDING + 0.2), Inches(1.9),
        Inches(CONTENT_WIDTH - 0.4), Inches(2.8)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = 'Courier New'
        p.font.color.rgb = hex_to_rgb(PALETTE['success'])
        p.space_after = Pt(4)

    return slide

def create_action_slide(prs, title, actions, section="", page_num=0):
    """9. Action Items Slide (액션 아이템)"""
    slide = add_slide(prs)
    set_background(slide, 'bg_primary')

    if section:
        add_section_badge(slide, section)
    if page_num:
        add_page_number(slide, page_num)

    # Title
    add_text_box(slide, title,
                 left=PADDING, top=1.0, width=CONTENT_WIDTH, height=0.6,
                 style='title', color_key='text_primary')

    # Action Cards
    card_height = 0.9
    start_y = 1.8

    for i, action in enumerate(actions):
        y = start_y + (i * (card_height + 0.15))

        # Card
        add_shape_with_fill(slide, PADDING, y, CONTENT_WIDTH, card_height, 'bg_card')

        # Number Circle
        add_text_box(slide, str(i + 1),
                     left=PADDING + 0.2, top=y + 0.2, width=0.5, height=0.5,
                     style='subtitle', color_key='accent', align='center', valign='middle')

        # Action Text
        add_text_box(slide, action,
                     left=PADDING + 0.9, top=y + 0.25, width=CONTENT_WIDTH - 1.3, height=0.5,
                     style='body', color_key='text_primary', valign='middle')

    return slide

def create_closing_slide(prs, title, subtitle="", contact=""):
    """10. Closing Slide (마무리)"""
    slide = add_slide(prs)
    set_background(slide, 'accent')

    # Title
    add_text_box(slide, title,
                 left=PADDING, top=2.0, width=CONTENT_WIDTH, height=0.8,
                 style='hero', color_key='text_primary', align='center', valign='middle')

    # Subtitle
    if subtitle:
        add_text_box(slide, subtitle,
                     left=PADDING, top=2.9, width=CONTENT_WIDTH, height=0.5,
                     style='subtitle', color_key='text_primary', align='center')

    # Contact
    if contact:
        add_text_box(slide, contact,
                     left=PADDING, top=4.2, width=CONTENT_WIDTH, height=0.4,
                     style='body', color_key='text_primary', align='center')

    return slide

# ============================================================
# 메인: Claude Code 2.0 세미나 PPT 생성
# ============================================================

def generate_claude_code_seminar():
    """Claude Code 2.0 세미나 PPT 생성"""
    prs = create_presentation()

    # ========================================
    # Slide 1: Cover
    # ========================================
    create_cover_slide(prs,
        title="Claude Code 2.0\n완벽 활용 가이드",
        subtitle="AI 시대, 개발 생산성을 폭발시키는 방법",
        presenter="Tech Team",
        date="2025.01")

    # ========================================
    # Slide 2: Contents
    # ========================================
    create_contents_slide(prs, [
        {'title': 'Why Claude Code', 'page': 3},
        {'title': 'Context Engineering', 'page': 5},
        {'title': '병렬 Sub-agents', 'page': 9},
        {'title': 'Skills & Commands', 'page': 14},
        {'title': 'MCP 통합', 'page': 18},
        {'title': '실전 & 마인드셋', 'page': 21},
        {'title': 'Closing', 'page': 24},
    ])

    # ========================================
    # Part 1: Why Claude Code (3-4)
    # ========================================
    create_section_divider(prs, 1, "Why Claude Code",
                          "AI 코딩 어시스턴트의 새로운 기준")

    create_statistics_slide(prs,
        title="개발 속도 2-3배 향상의 비밀",
        metrics=[
            {'value': '2-3x', 'label': '개발 속도'},
            {'value': '200K', 'label': 'Context Window'},
            {'value': '4.5', 'label': 'Opus Version'},
        ],
        section="WHY", page_num=4)

    # ========================================
    # Part 2: Context Engineering (5-8)
    # ========================================
    create_section_divider(prs, 2, "Context Engineering",
                          "200K 토큰 공간을 효과적으로 관리하는 방법")

    create_content_slide(prs,
        headline="Context Window = 200K 토큰 공간",
        bullets=[
            "시스템 프롬프트, 대화 이력, 파일 내용 등 포함",
            "토큰 ≈ 단어의 3/4 (영어 기준)",
            "한글은 토큰 소모가 더 큼",
            "200K ≈ 소설 2-3권 분량",
        ],
        section="CONTEXT", page_num=6, visual_placeholder=True)

    create_content_slide(prs,
        headline="Context Rot - 성능 저하의 원인",
        bullets=[
            "Context가 80% 이상 차면 성능 급격히 저하",
            "오래된 정보가 혼란을 유발",
            "불필요한 파일 읽기로 공간 낭비",
            "해결책: /compact 명령어로 정리",
        ],
        section="CONTEXT", page_num=7)

    create_code_slide(prs,
        title="Status Line으로 실시간 모니터링",
        code_lines=[
            '# .claude/settings.json',
            '{',
            '  "statusLineEnabled": true,',
            '  "statusLineTemplate":',
            '    "[{model}] Tokens: {tokens}/{max} ({percent}%)"',
            '}',
        ],
        section="CONTEXT", page_num=8)

    # ========================================
    # Part 3: 병렬 Sub-agents (9-13)
    # ========================================
    create_section_divider(prs, 3, "병렬 Sub-agents",
                          "독립 Context의 분신들로 시간을 단축")

    create_content_slide(prs,
        headline="Sub-agent = 독립 Context의 분신들",
        bullets=[
            "메인 Agent와 별도의 Context 유지",
            "병렬 실행으로 작업 시간 단축",
            "각자 다른 파일/작업 동시 처리",
            "결과만 메인 Agent로 전달",
        ],
        section="AGENTS", page_num=10, visual_placeholder=True)

    create_comparison_slide(prs,
        title="순차 vs 병렬: 시간 62% 단축",
        before_items=[
            "파일 1 분석 → 완료",
            "파일 2 분석 → 완료",
            "파일 3 분석 → 완료",
            "총 소요시간: 30분",
        ],
        after_items=[
            "파일 1, 2, 3 동시 분석",
            "결과 통합",
            "총 소요시간: 12분",
            "62% 시간 절약!",
        ],
        section="AGENTS", page_num=11)

    create_content_slide(prs,
        headline="비용 폭발 주의! 49개 병렬의 교훈",
        bullets=[
            "병렬 Agent 수 = 비용 배수",
            "권장: 2-4개가 Sweet Spot",
            "5개 이상은 수익 체감",
            "모델 믹싱으로 비용 최적화",
        ],
        section="AGENTS", page_num=12)

    create_code_slide(prs,
        title="실전 프롬프트 패턴",
        code_lines=[
            '# 병렬 분석 요청',
            '"다음 3개 파일을 병렬로 분석해줘:',
            '- src/auth.ts',
            '- src/api.ts',
            '- src/db.ts',
            '각각 보안 취약점을 찾아줘"',
        ],
        section="AGENTS", page_num=13)

    # ========================================
    # Part 4: Skills & Commands (14-17)
    # ========================================
    create_section_divider(prs, 4, "Skills & Commands",
                          "개인 노하우를 팀 자산으로 변환")

    create_content_slide(prs,
        headline="Skills = 재사용 가능한 지식 모듈",
        bullets=[
            "SKILL.md 파일에 워크플로우 정의",
            ".claude/skills/ 폴더에 저장",
            "Git으로 팀 공유 가능",
            "복잡한 작업을 명령어 하나로!",
        ],
        section="SKILLS", page_num=15, visual_placeholder=True)

    create_code_slide(prs,
        title="SKILL.md 작성 예시",
        code_lines=[
            '---',
            'name: code-review',
            'description: 코드 리뷰 자동화',
            '---',
            '',
            '# Code Review Skill',
            '1. 파일 변경사항 분석',
            '2. 코딩 컨벤션 체크',
            '3. 보안 취약점 스캔',
        ],
        section="SKILLS", page_num=16)

    create_content_slide(prs,
        headline="Custom Commands로 워크플로우 자동화",
        bullets=[
            "/feature - 새 기능 개발 시작",
            "/review - 코드 리뷰 실행",
            "/test - 테스트 작성 및 실행",
            "/deploy - 배포 프로세스 시작",
        ],
        section="SKILLS", page_num=17)

    # ========================================
    # Part 5: MCP 통합 (18-20)
    # ========================================
    create_section_divider(prs, 5, "MCP 통합",
                          "AI의 USB-C 포트 - 무한 확장 가능")

    create_content_slide(prs,
        headline="MCP = Model Context Protocol",
        bullets=[
            "AI와 외부 시스템 연결 표준",
            "데이터베이스, API, 파일시스템 등",
            "플러그인처럼 기능 확장",
            "한 번 설정, 계속 사용",
        ],
        section="MCP", page_num=19, visual_placeholder=True)

    create_statistics_slide(prs,
        title="4대 MCP 서버 활용법",
        metrics=[
            {'value': 'Context7', 'label': '문서 검색'},
            {'value': 'Playwright', 'label': '브라우저'},
            {'value': 'DB', 'label': '데이터베이스'},
            {'value': 'API', 'label': '외부 연동'},
        ],
        section="MCP", page_num=20)

    # ========================================
    # Part 6: 실전 & 마인드셋 (21-23)
    # ========================================
    create_section_divider(prs, 6, "실전 & 마인드셋",
                          "AI-Native 개발자로 성장하기")

    create_content_slide(prs,
        headline="새 기능 개발 A to Z",
        bullets=[
            "1. 요구사항 정리 (프롬프트 작성)",
            "2. 설계 리뷰 (AI와 토론)",
            "3. 코드 생성 (병렬 작업)",
            "4. 테스트 & 리팩토링",
            "5. 코드 리뷰 & 배포",
        ],
        section="MINDSET", page_num=22)

    create_comparison_slide(prs,
        title="구현 → 설계에 시간 투자",
        before_items=[
            "구현 70%",
            "설계 15%",
            "테스트 15%",
            "야근이 일상...",
        ],
        after_items=[
            "설계 50%",
            "구현 20%",
            "검증 30%",
            "워라밸 확보!",
        ],
        section="MINDSET", page_num=23)

    # ========================================
    # Closing (24-25)
    # ========================================
    create_action_slide(prs,
        title="내일부터 바로 할 수 있는 3가지",
        actions=[
            "Day 1: Status Line 설정하고 Context 모니터링 시작",
            "Day 2: 반복 작업을 Skill로 만들어보기",
            "Day 3: 병렬 Agent로 코드 리뷰 자동화",
        ],
        section="ACTION", page_num=24)

    create_closing_slide(prs,
        title="Thank You",
        subtitle="Questions & Discussion",
        contact="AI가 코드를 짜는 시대,\n우리는 더 나은 설계자가 됩니다.")

    # ========================================
    # 파일 저장
    # ========================================
    output_path = os.path.join(os.path.dirname(__file__),
                               "Claude_Code_2.0_Seminar_v2.pptx")
    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    print(f"📊 총 슬라이드: {len(prs.slides)}장")

    return output_path

if __name__ == "__main__":
    generate_claude_code_seminar()
